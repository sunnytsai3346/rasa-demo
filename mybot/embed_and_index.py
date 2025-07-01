# Install requirements first:
# pip install sentence-transformers faiss-cpu PyMuPDF ollama
# One-time PDF load and FAISS save

import glob
import json
import os
from pathlib import Path
from typing import Dict, List
import faiss
import fitz  # PyMuPDF
import numpy as np
from sentence_transformers import SentenceTransformer
import uuid

import pdfplumber
from docx import Document
from pptx import Presentation
import pandas as pd
from bs4 import BeautifulSoup
from sentence_transformers import SentenceTransformer
import faiss
import numpy as np

from bs4 import BeautifulSoup
from docx import Document
from pptx import Presentation
import pandas as pd
import pdfplumber

PDF_DIR = os.path.join(os.path.dirname(__file__), "Actions","DATA","PDF")
DOC_FOLDER = os.path.join(os.path.dirname(__file__), "Actions","DATA","training_materials")
EMBED_MODEL = "intfloat/e5-large-v2"
INDEX_PATH = os.path.join(os.path.dirname(__file__), "Actions","DATA","FAISS_index")
CHUNK_SIZE = 300
CHUNK_OVERLAP = 50
METADATA_STORE = "vector_meta.pkl"

# 1. Load and chunk 
# def extract_text_from_file(filepath: str) -> str:
#     ext = Path(filepath).suffix.lower()
    
#     if ext == ".pdf":
#         with pdfplumber.open(filepath) as pdf:
#             return "\n".join(page.extract_text() or "" for page in pdf.pages)

#     elif ext == ".docx":
#         doc = Document(filepath)
#         return "\n".join(para.text for para in doc.paragraphs)

#     elif ext == ".pptx":
#         prs = Presentation(filepath)
#         return "\n".join(shape.text for slide in prs.slides for shape in slide.shapes if hasattr(shape, "text"))

#     elif ext == ".xlsx":
#         xls = pd.read_excel(filepath, sheet_name=None)
#         return "\n\n".join(f"Sheet: {sheet}\n" + df.astype(str).to_csv(index=False) for sheet, df in xls.items())

#     elif ext in [".html", ".htm"]:
#         with open(filepath, "r", encoding="utf-8") as f:
#             soup = BeautifulSoup(f, "lxml")
#             return soup.get_text(separator="\n")

#     return ""
def chunk_text(text, chunk_size=CHUNK_SIZE, overlap=CHUNK_OVERLAP):
    chunks = []
    for i in range(0, len(text), chunk_size - overlap):
        chunks.append(text[i:i+chunk_size])
    return chunks

def load_pdf_text(path):
    with pdfplumber.open(path) as pdf:
        return "\n".join(page.extract_text() or "" for page in pdf.pages)

def load_docx_text(path):
    doc = Document(path)
    content = []
    current_heading = ""
    for para in doc.paragraphs:
        if para.style.name.startswith("Heading"):
            current_heading = para.text
        elif para.text.strip():
            chunk = f"{current_heading}\n{para.text.strip()}"
            content.append(chunk)
    return "\n".join(content)

def load_pptx_text(path):
    prs = Presentation(path)
    slides = []
    for i, slide in enumerate(prs.slides):
        slide_text = "\n".join(shape.text for shape in slide.shapes if hasattr(shape, "text"))
        slides.append(f"Slide {i+1}:\n{slide_text}")
    return "\n".join(slides)

def load_excel_text(path):
    df_texts = []
    xls = pd.ExcelFile(path)
    for sheet in xls.sheet_names:
        df = xls.parse(sheet)
        df_texts.append(f"Sheet: {sheet}\n" + df.astype(str).fillna("").to_string(index=False))
    return "\n\n".join(df_texts)

def load_html_text(path):
    with open(path, encoding="utf-8") as f:
        soup = BeautifulSoup(f, "lxml")
        sections = []
        for tag in soup.find_all(["h1", "h2", "h3", "p", "li"]):
            sections.append(tag.get_text(strip=True))
        return "\n".join(sections)
    
def load_all_documents():    
    chunks = []
    loaders = {
        ".pdf": load_pdf_text,
        ".docx": load_docx_text,
        ".pptx": load_pptx_text,
        ".xlsx": load_excel_text,
        ".xls": load_excel_text,
        ".html": load_html_text,
    }
    for ext, loader in loaders.items():
        for file in glob.glob(f"{DOC_FOLDER}/**/*{ext}", recursive=True):
            print(f"Processing: {file}")            
            try:
                text = loader(file)
                for chunk in chunk_text(text):
                    chunks.append({"text": chunk, 
                                   "source": file,  # flattened for easy access later,
                                   "meta": {"file": os.path.basename(file)}})
            except Exception as e:
                print(f"Failed to parse {file}: {e}")
    return chunks
    # for root, dirs, files in os.walk(folder_path):
    #     for file in files:
    #         ext = Path(file).suffix.lower()
    #         if ext in supported_exts:
    #             filepath = os.path.join(root, file)
    #             print(f"Loading: {filepath}")
    #             text = extract_text_from_file(filepath)
    #             for i in range(0, len(text), CHUNK_SIZE):
    #                 chunk = text[i:i+CHUNK_SIZE].strip()
    #                 if chunk:
    #                     chunks.append({
    #                         "text": chunk,
    #                         "meta": {"file": file}
    #                     })    
    #     # if file.endswith(".pdf"):
    #     #     doc = fitz.open(os.path.join(pdf_dir, file))
    #     #     for page in doc:
    #     #         text = page.get_text()
    #     #         if len(text.strip()) > 100:
    #     #             chunks.append({"text": text,
    #     #                            "source": file,  # flattened for easy access later,
    #     #                            "meta": {"file": file}})
    # return chunks

# 2. Embed chunks
def embed_chunks(chunks: List[Dict], model: SentenceTransformer):
    texts = ["query: " + c["text"] for c in chunks]
    embeddings = model.encode(texts, convert_to_numpy=True, normalize_embeddings=True)
    return embeddings

# 3. Save to Faiss

def save_faiss_index(chunks: List[Dict], embeddings: np.ndarray):
    print('save_faiss_index')
    dim = embeddings.shape[1]
    index = faiss.IndexFlatIP(dim)
    index.add(embeddings)
    faiss.write_index(index, os.path.join(INDEX_PATH, "index.faiss"))

    # Save metadata
    with open(os.path.join(INDEX_PATH, METADATA_STORE), "w", encoding="utf-8") as f:
        json.dump(chunks, f)


# ---- To prepare the DB (run once):
if __name__ == "__main__":
    print("Loading and embedding PDFs...")
    model = SentenceTransformer(EMBED_MODEL)
    chunks = load_all_documents()
    embeddings = embed_chunks(chunks, model)
    save_faiss_index(chunks, embeddings)
    print("FAISS index saved.")

 
